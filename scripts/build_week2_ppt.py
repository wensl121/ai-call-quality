"""生成周分享 2 PPT（基于周分享 1 模板）。

用法：python scripts/build_week2_ppt.py
输出：C:\\Users\\Z\\Desktop\\周分享\\周分享2.pptx
"""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Emu, Pt

WEEK1 = Path(r"C:\Users\Z\Desktop\周分享\周分享1.pptx")
OUT = Path(r"C:\Users\Z\Desktop\周分享\周分享2.pptx")
SHOTS = Path(r"C:\Users\Z\Desktop\周分享\周分享2")

PRIMARY = RGBColor(0xC0, 0x39, 0x2B)        # 招商红
ACCENT = RGBColor(0x1F, 0x49, 0x7D)
DARK = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0xEE, 0xEE, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _delete_all_slides(prs: Presentation) -> None:
    """保留 master / layouts，把 slides 全部清空。"""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"]
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _layout(prs: Presentation, name: str):
    for lo in prs.slide_layouts:
        if lo.name == name:
            return lo
    return prs.slide_layouts[1]  # fallback "Title and Content"


def _add_text_box(
    slide,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    text: str,
    *,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = DARK,
    align=PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = "微软雅黑"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def _add_bullets(
    slide,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    items: list[str],
    *,
    font_size: int = 13,
    color: RGBColor = DARK,
    bullet_color: RGBColor = PRIMARY,
):
    tb = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(4)

        bul = para.add_run()
        bul.text = "• "
        bul.font.name = "微软雅黑"
        bul.font.size = Pt(font_size)
        bul.font.bold = True
        bul.font.color.rgb = bullet_color

        run = para.add_run()
        run.text = item
        run.font.name = "微软雅黑"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return tb


def _add_section_header(slide, number: str, title: str):
    """页面左上角的 "01 标题" 风格 header。"""
    # 数字
    num = slide.shapes.add_textbox(Cm(0.8), Cm(0.4), Cm(2), Cm(1.5))
    tf = num.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = number
    r.font.name = "Arial"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    # 标题
    tit = slide.shapes.add_textbox(Cm(2.6), Cm(0.7), Cm(28), Cm(1.2))
    tf = tit.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "微软雅黑"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = DARK
    # 红色下划线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0.8), Cm(2.0), Cm(1.5), Cm(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()


def _add_code_block(slide, left_cm, top_cm, width_cm, height_cm, code: str, font_size: int = 10):
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Cm(0.3)
    tf.margin_right = Cm(0.3)
    tf.margin_top = Cm(0.2)
    tf.margin_bottom = Cm(0.2)
    tf.word_wrap = True

    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(font_size)
        r.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)


def _add_picture_safe(slide, path: Path, left_cm, top_cm, width_cm=None, height_cm=None):
    if not path.exists():
        return None
    kwargs = {"left": Cm(left_cm), "top": Cm(top_cm)}
    if width_cm:
        kwargs["width"] = Cm(width_cm)
    if height_cm:
        kwargs["height"] = Cm(height_cm)
    return slide.shapes.add_picture(str(path), **kwargs)


def _add_footer(slide, page_no: int, total: int):
    """右下角页码。"""
    tb = slide.shapes.add_textbox(Cm(30), Cm(18.7), Cm(3), Cm(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page_no} / {total}"
    r.font.name = "Arial"
    r.font.size = Pt(9)
    r.font.color.rgb = GREY


# ============================================================
# Main builder
# ============================================================

def build():
    prs = Presentation(str(WEEK1))
    _delete_all_slides(prs)

    title_layout = _layout(prs, "Title and Content")
    blank_layout = _layout(prs, "仅标题") if any(l.name == "仅标题" for l in prs.slide_layouts) else title_layout

    # Slide width/height in cm (16:9)
    sw = prs.slide_width / 360000  # EMU → cm
    sh = prs.slide_height / 360000

    pages = []  # (build_func, ...)

    # ===== 1. 封面 =====
    def slide_1():
        s = prs.slides.add_slide(title_layout)
        _add_text_box(s, 1.5, 5.5, 30, 2,
                      "Agent 技术框架与产品实践分享 02",
                      font_size=32, bold=True, color=DARK)
        _add_text_box(s, 1.5, 8.5, 30, 1.5,
                      "—— LangChain & LangGraph + AI 通话质检系统实战",
                      font_size=18, color=GREY)
        # 红色装饰条
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.5), Cm(11), Cm(2.5), Cm(0.15))
        bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background()
        _add_text_box(s, 1.5, 16.5, 15, 0.8,
                      "MORE  VALUE", font_size=13, bold=True, color=PRIMARY)
        _add_text_box(s, 1.5, 17.3, 15, 0.8,
                      "FOR INVESTORS", font_size=13, bold=True, color=PRIMARY)
        _add_text_box(s, 23, 16.8, 10, 0.7, "分享人：温舒麟", font_size=12, color=GREY, align=PP_ALIGN.RIGHT)
        _add_text_box(s, 23, 17.7, 10, 0.7, "部门：信息技术部 - 开发三组", font_size=12, color=GREY, align=PP_ALIGN.RIGHT)
        return s

    # ===== 2. 目录 =====
    def slide_2():
        s = prs.slides.add_slide(blank_layout)
        _add_text_box(s, 1.5, 1.0, 25, 1.5, "本周分享 · 三大模块",
                      font_size=24, bold=True, color=DARK)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.5), Cm(2.6), Cm(2.5), Cm(0.12))
        bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background()

        # 三大模块卡片
        items = [
            ("01", "LangChain & LangGraph 概念", "为什么需要图编排 / 核心组件 / Middleware"),
            ("02", "项目实战：AI 通话质检", "从 Dify 迁移到 LangGraph 的全过程，5 项核心能力"),
            ("03", "LangSmith 工程化", "Tracing · Evals · Studio · Platform 部署"),
        ]
        top = 4.5
        for num, title, sub in items:
            _add_text_box(s, 1.5, top, 2, 1.5, num, font_size=36, bold=True, color=PRIMARY)
            _add_text_box(s, 4.5, top + 0.2, 25, 1, title, font_size=18, bold=True, color=DARK)
            _add_text_box(s, 4.5, top + 1.2, 25, 0.8, sub, font_size=12, color=GREY)
            top += 3.5
        return s

    # ===== 3. 01-1 LangChain vs LangGraph (上) =====
    def slide_3():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "01", "LangChain Agent vs LangGraph：何时用哪个？")
        _add_text_box(s, 1, 3, 30, 0.8,
                      "LangChain 1.0 把这两条路线明确分开",
                      font_size=14, color=GREY)
        # 对比表
        rows = [
            ("",                   "LangChain Agent（create_agent）",      "LangGraph"),
            ("适用",               "单一 Agent（一个模型 + 工具循环）",     "多步骤工作流 / 复杂分支与循环"),
            ("控制粒度",           "中等（middleware 钩子）",                "高（每条边、每个节点都能定制）"),
            ("典型场景",           "ReAct、客服 chatbot、tool-using agent",  "审批流、质检、研究流水线、HITL"),
            ("状态",               "消息历史为主",                           "TypedDict + 自定义 reducer"),
            ("可观测",             "LangSmith 自动",                         "LangSmith 自动 + Studio 可视化调试"),
        ]
        from pptx.util import Cm as C
        n_cols = 3
        col_widths = [4.5, 12, 12.5]
        left = 1.5
        top = 4.4
        row_h = 1.4
        for r_idx, row in enumerate(rows):
            x = left
            for c_idx, cell in enumerate(row):
                rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, C(x), C(top + r_idx * row_h),
                                          C(col_widths[c_idx]), C(row_h))
                if r_idx == 0:
                    rect.fill.solid(); rect.fill.fore_color.rgb = PRIMARY
                    rect.line.fill.background()
                    text_color = WHITE; bold = True
                else:
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = WHITE if r_idx % 2 else LIGHT_GREY
                    rect.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                    text_color = DARK; bold = (c_idx == 0)
                tf = rect.text_frame
                tf.margin_left = C(0.2); tf.margin_right = C(0.2)
                tf.margin_top = C(0.15); tf.margin_bottom = C(0.15)
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = cell
                run.font.name = "微软雅黑"
                run.font.size = Pt(11)
                run.font.bold = bold
                run.font.color.rgb = text_color
                x += col_widths[c_idx]

        _add_text_box(s, 1.5, 14.2, 30, 0.8,
                      "👉 我们的质检系统：5 节点 + 子图循环 + 并行打分 → 选 LangGraph",
                      font_size=13, bold=True, color=ACCENT)
        return s

    # ===== 4. 01-1 (下) Dify 流程图说明 =====
    def slide_4():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "01", "为什么从 Dify 迁到 LangGraph？")
        _add_picture_safe(s, SHOTS / "dify流程图.jpg", 1.5, 3.0, width_cm=22)
        _add_text_box(s, 24.5, 3.3, 9, 0.8, "原 Dify 流程", font_size=14, bold=True, color=DARK)
        _add_bullets(s, 24.5, 4.2, 9, 6, [
            "可视化拼装 ✓",
            "节点逻辑黑盒",
            "状态难管理",
            "调试只能看节点输出",
            "测试无回归保障",
            "成本不可见",
        ], font_size=11)
        _add_text_box(s, 24.5, 11.5, 9, 0.8, "迁到 LangGraph 后", font_size=14, bold=True, color=PRIMARY)
        _add_bullets(s, 24.5, 12.4, 9, 5, [
            "代码化 / 可 Code Review",
            "TypedDict 强类型 state",
            "trace 全链路可视化",
            "Pytest + Evals 双层回归",
            "每次调用 token 成本可统计",
        ], font_size=11)
        return s

    # ===== 5. 01-2 Graph Runtime 概念 =====
    def slide_5():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "01", "Graph Runtime 核心概念")
        _add_text_box(s, 1.5, 3, 30, 0.8, "把 Agent 流程画成有向图，框架管 state、调度、并发、循环",
                      font_size=14, color=GREY)
        # 4 个概念卡片
        cards = [
            ("Node 节点",     "一个函数：state → state 增量\n可纯计算，也可调 LLM/工具", PRIMARY),
            ("Edge 边",       "节点间的连接\n固定边 / 条件边 / Send 动态边",          ACCENT),
            ("State 状态",    "TypedDict + Reducer\n跨节点共享，自动合并",            RGBColor(0x2E, 0x7D, 0x32)),
            ("Subgraph 子图", "把一组节点封成单元\n父子图自动共享 state 键",          RGBColor(0xEF, 0x6C, 0x00)),
        ]
        x0 = 1.5
        y0 = 4.5
        cw = 7.5
        ch = 5.5
        gap = 0.5
        for i, (title, desc, color) in enumerate(cards):
            x = x0 + i * (cw + gap)
            top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y0), Cm(cw), Cm(0.5))
            top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()
            body = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y0 + 0.5), Cm(cw), Cm(ch - 0.5))
            body.fill.solid(); body.fill.fore_color.rgb = WHITE
            body.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            _add_text_box(s, x + 0.4, y0 + 0.8, cw - 0.8, 1, title, font_size=15, bold=True, color=color)
            _add_text_box(s, x + 0.4, y0 + 2.2, cw - 0.8, 3, desc, font_size=12, color=DARK)

        _add_text_box(s, 1.5, 11, 30, 0.6, "执行模型", font_size=14, bold=True, color=DARK)
        _add_bullets(s, 1.5, 11.8, 30, 4, [
            "运行时按 BFS 推进：当前层所有节点完成 → 触发下一层",
            "状态更新通过 reducer 合并，支持并行写入同一字段",
            "条件边返回 Send(node, payload) → 动态 fan-out 多个并行实例",
            "compile() 后是不可变对象，多次 invoke 共享同一 graph",
        ], font_size=12)
        return s

    # ===== 6. 01-3 State / Reducer / Send / Subgraph =====
    def slide_6():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "01", "State + Reducer + Send：并行编排的三件套")
        _add_text_box(s, 1.5, 3, 16, 0.6, "1. TypedDict + Annotated reducer", font_size=13, bold=True, color=PRIMARY)
        _add_code_block(s, 1.5, 3.7, 16, 4, """class GraphState(TypedDict, total=False):
    conversation: str
    deductions: Annotated[
        list[dict],
        list_add_or_reset    # 自定义 reducer
    ]
    fatal_triggers: Annotated[
        list[dict], list_add_or_reset
    ]""", font_size=11)

        _add_text_box(s, 1.5, 8.2, 16, 0.6, "2. 自定义 reducer（支持 RESET 信号）", font_size=13, bold=True, color=PRIMARY)
        _add_code_block(s, 1.5, 8.9, 16, 4, """def list_add_or_reset(left, right):
    if right is RESET:
        return []                    # 清空
    return (left or []) + (right or [])

# 重打分循环开头：
return {"deductions": RESET, ...}""", font_size=11)

        _add_text_box(s, 18.3, 3, 14, 0.6, "3. Send：动态 fan-out", font_size=13, bold=True, color=ACCENT)
        _add_code_block(s, 18.3, 3.7, 14, 6, """def fan_out(state):
    return [
        Send("rule_scorer", {
            "rule": rule,
            "conversation": state["conversation"],
        })
        for rule in state["rules_json"]
    ]

# 5 条规则 → 5 个 rule_scorer
# 并行跑，结果通过 reducer 累加""", font_size=11)
        _add_text_box(s, 18.3, 10.2, 14, 0.6, "4. Subgraph：模块化封装", font_size=13, bold=True, color=ACCENT)
        _add_code_block(s, 18.3, 10.9, 14, 4, """sub = build_scoring_subgraph()
parent.add_node(
    "scoring_with_audit", sub
)
# 父子图共享同名 state 键""", font_size=11)
        return s

    # ===== 7. 01-4 Middleware & Context Engineering =====
    def slide_7():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "01", "Middleware & Context Engineering")
        _add_text_box(s, 1.5, 3, 15, 0.8, "Middleware（横切关注点）", font_size=16, bold=True, color=PRIMARY)
        _add_text_box(s, 1.5, 4.2, 15, 0.6, "在 LLM 调用前后插入逻辑，业务节点零侵入", font_size=12, color=GREY)
        _add_bullets(s, 1.5, 5.3, 15, 6, [
            "before_model：注入上下文 / PII 脱敏 / 权限校验",
            "after_model：日志 / 成本统计 / 输出审查",
            "wrap_model_call：完全包装（重试 / 缓存 / 故障切换）",
            "在 LangGraph 里用 Runnable.with_listeners 或自建 helper",
        ], font_size=12)

        _add_text_box(s, 17, 3, 15, 0.8, "Context Engineering", font_size=16, bold=True, color=ACCENT)
        _add_text_box(s, 17, 4.2, 15, 0.6, "重点不再是\"写提示词\"，而是\"设计上下文系统\"", font_size=12, color=GREY)
        _add_bullets(s, 17, 5.3, 15, 6, [
            "Prompt Engineering：单条 prompt 怎么写好",
            "Context Engineering：怎么动态拼上下文",
            "  · 历史压缩 / 摘要",
            "  · 知识检索（RAG）",
            "  · few-shot 示例注入",
            "  · 结构化输出 schema",
            "  · 工具描述与权限",
        ], font_size=12)

        _add_text_box(s, 1.5, 14.5, 30, 1,
                      "👉 本周项目实战中的 with_structured_output / KB 检索 / 审核反馈，都是 Context Engineering",
                      font_size=12, bold=True, color=ACCENT)
        return s

    # ===== 8. 02-1 业务背景 =====
    def slide_8():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "02", "项目背景：金融客服通话智能质检")
        _add_text_box(s, 1.5, 3, 14, 0.8, "业务", font_size=14, bold=True, color=PRIMARY)
        _add_bullets(s, 1.5, 4, 14, 5, [
            "对客服-客户通话录音做合规、话术、态度多维度评分",
            "原系统用 Dify 拼装：LLM 解析 + 规则 JSON + 扣分输出",
            "目标：迁到 LangGraph，提高可解释性与可扩展性",
        ], font_size=12)

        _add_text_box(s, 1.5, 9.3, 14, 0.8, "评分维度", font_size=14, bold=True, color=PRIMARY)
        _add_bullets(s, 1.5, 10.3, 14, 5, [
            "T1 合规风险事项（fatal）",
            "T2 服务态度（fatal）",
            "T3 服务标准性及服务技巧",
            "T4 业务水平",
        ], font_size=12)

        _add_text_box(s, 17, 3, 15, 0.8, "技术目标", font_size=14, bold=True, color=ACCENT)
        _add_bullets(s, 17, 4, 15, 7, [
            "保留原始对话文本，不清洗不拆分",
            "问题提取 + KB 检索：客户问题 → 知识库参考答案",
            "规则打分：循环最多 3 次",
            "审核：3 次未通过 → 标记人工复核",
            "输出：caps / deductions / hot_words / cost_summary",
            "可扩展性：未来加新规则不改框架代码",
        ], font_size=12)
        return s

    # ===== 9. 02-2 架构演进 =====
    def slide_9():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "02", "架构演进：从 6 节点线性 到 5 节点 + 子图")
        _add_text_box(s, 1.5, 3, 15, 0.7, "v1 初始版本", font_size=14, bold=True, color=GREY)
        _add_picture_safe(s, SHOTS / "初始图.png", 1.5, 3.8, width_cm=15)
        _add_text_box(s, 17, 3, 15, 0.7, "v2 当前版本（含子图 + 并行）", font_size=14, bold=True, color=PRIMARY)
        _add_picture_safe(s, SHOTS / "节点图.png", 17, 3.8, width_cm=15)

        _add_text_box(s, 1.5, 14.3, 30, 0.8, "演进收益", font_size=13, bold=True, color=DARK)
        _add_bullets(s, 1.5, 15.1, 30, 4, [
            "scoring_loop ⇄ auditor 抽成子图 → 主图只看 5 步流水线",
            "scoring_loop 拆成 dispatcher + 多个 rule_scorer + extraction → Send 并行",
            "Studio 中子图可双击展开，调试更直观",
        ], font_size=12)
        return s

    # ===== 10. 02-3 结构化输出 =====
    def slide_10():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "02", "核心①：with_structured_output 强类型 LLM 输出")
        _add_text_box(s, 1.5, 3, 30, 0.8, "痛点：LLM 直接返回 JSON 字符串，json.loads 偶尔崩；偷懒返回空数组",
                      font_size=12, color=GREY)
        _add_text_box(s, 1.5, 4, 15, 0.6, "Pydantic schema", font_size=13, bold=True, color=PRIMARY)
        _add_code_block(s, 1.5, 4.7, 15, 5.5, """class RuleScoringOutput(BaseModel):
    deductions: list[Deduction]
    fatal_triggers: list[FatalTrigger]

class Deduction(BaseModel):
    rule_id: str
    severity: Literal["fatal","nonfatal"]
    subtotal: int
    evidence: list[Evidence]""", font_size=11)
        _add_text_box(s, 17, 4, 15, 0.6, "节点调用", font_size=13, bold=True, color=PRIMARY)
        _add_code_block(s, 17, 4.7, 15, 5.5, """def rule_scorer(payload):
    parsed, usage = invoke_structured(
        RuleScoringOutput,
        [SystemMessage(_PROMPT),
         HumanMessage(payload_json)],
        node_name="rule_scorer",
    )
    return {"deductions": ..., "llm_usage": [usage]}""", font_size=11)
        _add_text_box(s, 1.5, 10.8, 30, 0.6, "实测对比（同一段对话）", font_size=13, bold=True, color=DARK)
        rows = [
            ("字段",          "raw JSON 解析",      "structured output"),
            ("hot_words",     "[]（模型偷懒）",     "[基金, 招商, 白酒, ...]"),
            ("business_words","[]",                 "[持仓查询]"),
            ("should_say",    "[]",                 "[\"您好，招商基金...\"]"),
        ]
        cw = [6, 12, 12]
        x = 1.5; y = 11.7
        for ri, row in enumerate(rows):
            cx = x
            for ci, cell in enumerate(row):
                rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(cx), Cm(y + ri * 0.95), Cm(cw[ci]), Cm(0.95))
                if ri == 0:
                    rect.fill.solid(); rect.fill.fore_color.rgb = PRIMARY; rect.line.fill.background()
                    color = WHITE; bold = True
                else:
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT_GREY
                    rect.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                    color = DARK; bold = (ci == 0)
                tf = rect.text_frame
                tf.margin_left = Cm(0.2); tf.margin_top = Cm(0.1); tf.margin_bottom = Cm(0.1)
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = cell
                run.font.name = "微软雅黑"
                run.font.size = Pt(11); run.font.bold = bold
                run.font.color.rgb = color
                cx += cw[ci]
        return s

    # ===== 11. 02-4 Send 并行 =====
    def slide_11():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "02", "核心②：Send 并行打分 + State Reducer")
        _add_text_box(s, 1.5, 3, 30, 0.8, "痛点：30 条规则塞一个 LLM 调用 → prompt 太长、漏判、延迟 O(n)",
                      font_size=12, color=GREY)
        _add_code_block(s, 1.5, 4, 19, 7, """# scoring_dispatcher: 重置 + fan-out
def scoring_dispatcher(state):
    return {"deductions": RESET, "fatal_triggers": RESET}

def fan_out(state):
    sends = [Send("rule_scorer", {"rule": r, ...})
             for r in state["rules_json"]]
    sends.append(Send("extraction_node", {...}))
    return sends

# rule_scorer × N 并行写 deductions（reducer 累加）
# extraction_node 并行抽 hot_words / should_say
# 所有 Send 完成 → 自动 fan-in 到 auditor""", font_size=11)
        _add_text_box(s, 21.5, 4, 11, 0.6, "性能", font_size=13, bold=True, color=PRIMARY)
        rows = [
            ("规则数",  "原方案",   "Send 并行"),
            ("1",       "10 s",     "12 s"),
            ("5",       "~10 s",    "12.5 s"),
            ("30 (估)", "30-60 s",  "~13-15 s"),
        ]
        cw = [3.5, 3.7, 3.8]
        x = 21.5; y = 4.8
        for ri, row in enumerate(rows):
            cx = x
            for ci, cell in enumerate(row):
                rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(cx), Cm(y + ri * 1.0), Cm(cw[ci]), Cm(1.0))
                if ri == 0:
                    rect.fill.solid(); rect.fill.fore_color.rgb = PRIMARY; rect.line.fill.background()
                    color = WHITE; bold = True
                else:
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT_GREY
                    rect.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                    color = DARK; bold = (ci == 0 or ri == 3)
                tf = rect.text_frame
                tf.margin_left = Cm(0.15); tf.margin_top = Cm(0.15)
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                run = p.add_run(); run.text = cell
                run.font.name = "微软雅黑"; run.font.size = Pt(11)
                run.font.bold = bold; run.font.color.rgb = color
                cx += cw[ci]

        _add_text_box(s, 21.5, 9.5, 11, 4,
                      "规则越多，并行收益越大；\n每个 scorer prompt 短了，\n准确度也提升（不会因\ncontext 太长丢规则）",
                      font_size=11, color=DARK)

        _add_text_box(s, 1.5, 12, 30, 0.6, "Reducer 解决的关键问题",
                      font_size=13, bold=True, color=ACCENT)
        _add_bullets(s, 1.5, 12.8, 30, 4, [
            "默认行为：多个并行节点写同一字段 → 后写覆盖前写（数据丢失）",
            "Annotated[list, list_add_or_reset]：自动 append；遇 RESET 标志清空",
            "审核未通过回到 dispatcher → RESET 让本轮从头累加，不污染历史",
        ], font_size=12)
        return s

    # ===== 12. 02-5 Subgraph =====
    def slide_12():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "02", "核心③：Subgraph 抽取打分-审核闭环")
        _add_text_box(s, 1.5, 3, 30, 0.8,
                      "把 scoring_dispatcher → rule_scorer × N → auditor 这一段封成子图",
                      font_size=13, color=GREY)

        _add_text_box(s, 1.5, 4, 15, 0.6, "主图（5 节点）", font_size=13, bold=True, color=PRIMARY)
        _add_code_block(s, 1.5, 4.7, 15, 6, """def build_graph():
    g = StateGraph(GraphState)
    g.add_node("input", input_node)
    g.add_node("question_extractor", ...)
    g.add_node("knowledge_retriever", ...)
    g.add_node(
        "scoring_with_audit",
        build_scoring_subgraph()    # ← 子图作节点
    )
    g.add_node("aggregator", aggregator)""", font_size=11)

        _add_text_box(s, 17, 4, 15, 0.6, "子图（4 节点 + 循环）", font_size=13, bold=True, color=ACCENT)
        _add_code_block(s, 17, 4.7, 15, 6, """def build_scoring_subgraph():
    g = StateGraph(GraphState)
    g.add_node("scoring_dispatcher", ...)
    g.add_node("rule_scorer", rule_scorer)
    g.add_node("extraction_node", ...)
    g.add_node("auditor", auditor)
    # 内部循环 + Send fan-out
    return g.compile()""", font_size=11)

        _add_text_box(s, 1.5, 11.5, 30, 0.6, "收益", font_size=13, bold=True, color=DARK)
        _add_bullets(s, 1.5, 12.3, 30, 6, [
            "主图清晰：从外面看就是一条直线流水线，循环细节封在内部",
            "可独立测试：build_scoring_subgraph() 可单独 compile + 单测",
            "可复用：同样的 \"draft + critique 循环\" 模式可用在其他地方",
            "Studio 可视化：子图节点可双击展开，调试更友好",
            "状态自动透传：父子图共享 GraphState，零 plumbing 代码",
        ], font_size=12)
        return s

    # ===== 13. 02-6 Middleware: PII =====
    def slide_13():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "02", "核心④：Middleware - PII 脱敏（合规刚需）")
        _add_text_box(s, 1.5, 3, 30, 0.8,
                      "金融通话含手机号 / 身份证 / 银行卡 / 地址，进 LLM 前必须脱敏",
                      font_size=12, color=GREY)

        _add_text_box(s, 1.5, 4, 15, 0.6, "middleware/pii.py", font_size=13, bold=True, color=PRIMARY)
        _add_code_block(s, 1.5, 4.7, 15, 6.5, """_PATTERNS = [
    (re.compile(r"\\b[\\w.+-]+@[\\w-]+\\.[\\w.-]+\\b"),
     "[EMAIL]"),
    (re.compile(r"\\b\\d{17}[\\dXx]\\b"),
     "[ID_CARD]"),
    (re.compile(r"(?<!\\d)1[3-9]\\d{9}(?!\\d)"),
     "[PHONE]"),
    (re.compile(r"[一-龥]{2,}(?:路|街|巷)\\d+号..."),
     "[ADDRESS]"),
]
def redact_pii(text):
    for p, r in _PATTERNS:
        text = p.sub(r, text)
    return text""", font_size=11)

        _add_text_box(s, 17, 4, 15, 0.6, "在 input_node 一次性应用", font_size=13, bold=True, color=ACCENT)
        _add_code_block(s, 17, 4.7, 15, 4, """def input_node(state):
    return {
        "conversation": redact_pii(
            state["conversation"]
        ),
        ...
    }""", font_size=11)

        _add_text_box(s, 17, 9.3, 15, 0.6, "效果", font_size=13, bold=True, color=ACCENT)
        _add_bullets(s, 17, 10, 15, 5, [
            "下游所有节点（含 LLM）只看 [PHONE] [ID_CARD]",
            "evidence.text 引用对话也是脱敏版",
            "对外输出报告天然合规",
        ], font_size=11)

        _add_text_box(s, 1.5, 14, 30, 0.8,
                      "👉 业务节点零侵入：input_node 加一行调用，全链路自动脱敏",
                      font_size=12, bold=True, color=ACCENT)
        return s

    # ===== 14. 02-7 成本追踪 =====
    def slide_14():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "02", "核心⑤：成本追踪（每个节点的 token 与花费）")
        _add_text_box(s, 1.5, 3, 30, 0.8,
                      "include_raw=True 拿 AIMessage.usage_metadata，list reducer 累加",
                      font_size=12, color=GREY)
        _add_picture_safe(s, SHOTS / "成本计算.png", 1.5, 4, width_cm=20)
        _add_text_box(s, 22.5, 4, 10, 0.7, "看到了什么", font_size=13, bold=True, color=PRIMARY)
        _add_bullets(s, 22.5, 5, 10, 7, [
            "rule_scorer 5 次并行调用",
            "question_extractor 跑了 2 次（隐藏成本）",
            "9 次 LLM 共 15K input + 1.1K output tokens",
            "单通成本 $0.0024（DeepSeek）",
            "可推算月成本上限",
        ], font_size=11)
        _add_text_box(s, 22.5, 12.3, 10, 2,
                      "👉 量化是优化的前提",
                      font_size=14, bold=True, color=ACCENT)
        return s

    # ===== 15. 03-1 Tracing =====
    def slide_15():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "03", "LangSmith Tracing：每次调用全链路可视化")
        _add_picture_safe(s, SHOTS / "追踪.png", 1.5, 3, width_cm=22)
        _add_text_box(s, 24.5, 3.3, 9, 0.7, "Trace 视图特点", font_size=13, bold=True, color=PRIMARY)
        _add_bullets(s, 24.5, 4.2, 9, 6, [
            "完整调用树：每个节点 / 每次 LLM",
            "Waterfall 时间线：哪个慢一目了然",
            "input / output / usage 全可看",
            "失败 trace 标红，可重放",
            "支持搜索、过滤、分组",
        ], font_size=11)
        _add_text_box(s, 24.5, 11, 9, 0.7, "如何接入", font_size=13, bold=True, color=ACCENT)
        _add_bullets(s, 24.5, 11.9, 9, 4, [
            ".env 配 4 个变量",
            "LangChain ChatOpenAI 自动上传",
            "无需写一行追踪代码",
        ], font_size=11)
        return s

    # ===== 16. 03-2 Evals + 黄金数据集 =====
    def slide_16():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "03", "Evals 回归：黄金数据集 + 自动评估器")
        _add_text_box(s, 1.5, 3, 30, 0.7,
                      "改 prompt 后\"感觉变好了\" → \"rule_match 从 75% → 92%\" 的转变",
                      font_size=12, color=GREY)
        _add_picture_safe(s, SHOTS / "黄金数据库和实验.png", 1.5, 3.9, width_cm=15)
        _add_picture_safe(s, SHOTS / "柱状图.png", 17, 3.9, width_cm=15)
        _add_text_box(s, 1.5, 12.3, 15, 0.7, "黄金数据集（4 类例子）", font_size=12, bold=True, color=PRIMARY)
        _add_bullets(s, 1.5, 13, 15, 4, [
            "clean_compliant：合规通话",
            "missing_greeting：缺话术",
            "fatal_guarantee：致命违规",
            "kb_mismatch：答非所问",
        ], font_size=11)
        _add_text_box(s, 17, 12.3, 15, 0.7, "4 个评估器（baseline 实测）", font_size=12, bold=True, color=ACCENT)
        _add_bullets(s, 17, 13, 15, 4, [
            "rule_match: 75%（Jaccard 命中）",
            "fatal_correctness: 75% (精确)",
            "score_in_range: 100%",
            "review_match: 100%",
        ], font_size=11)
        return s

    # ===== 17. 03-3 Studio + Platform 部署 =====
    def slide_17():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "03", "Studio 调试 + Platform 一键部署")
        _add_text_box(s, 1.5, 3, 15, 0.7, "本地 Studio（langgraph dev）", font_size=13, bold=True, color=PRIMARY)
        _add_picture_safe(s, SHOTS / "节点图.png", 1.5, 3.8, width_cm=15)
        _add_text_box(s, 17, 3, 15, 0.7, "Platform Deployment", font_size=13, bold=True, color=ACCENT)
        _add_picture_safe(s, SHOTS / "部署.png", 17, 3.8, width_cm=15)
        _add_text_box(s, 1.5, 14, 30, 0.6, "工作流", font_size=13, bold=True, color=DARK)
        _add_bullets(s, 1.5, 14.7, 30, 4, [
            "本地：langgraph dev → Studio 浏览器调试，节点图实时展开 / 重放",
            "线上：langgraph.json 配置好 → Push GitHub → LangSmith Platform 一键部署成 HTTPS API",
            "环境变量在 Platform UI 配（DEEPSEEK_API_KEY 等）",
        ], font_size=12)
        return s

    # ===== 18. 03-4 Claude Code 辅助 =====
    def slide_18():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "03", "彩蛋：整个项目用 Claude Code 协作完成")
        _add_picture_safe(s, SHOTS / "claude工作.png", 1.5, 3, width_cm=18)
        _add_text_box(s, 20.5, 3, 12, 0.7, "用 AI 写 AI 系统的体感", font_size=13, bold=True, color=PRIMARY)
        _add_bullets(s, 20.5, 4, 12, 8, [
            "需求 → 项目骨架：30 分钟",
            "增量重构（结构化输出 → Send 并行 → Subgraph）：每次 ~1 小时",
            "Middleware（PII + 成本）零侵入加上去",
            "Evals 黄金数据集 + 评估器代码全自动",
            "Bug 现场调试：贴报错就能定位",
        ], font_size=11)
        _add_text_box(s, 20.5, 12, 12, 1,
                      "👉 第 3 周分享主题：Claude Agent SDK & MCP",
                      font_size=12, bold=True, color=ACCENT)
        return s

    # ===== 19. 总结 =====
    def slide_19():
        s = prs.slides.add_slide(title_layout)
        _add_section_header(s, "结", "总结与下周方向")
        _add_text_box(s, 1.5, 3, 15, 0.8, "本周关键收获", font_size=15, bold=True, color=PRIMARY)
        _add_bullets(s, 1.5, 4.2, 15, 8, [
            "LangGraph 适合需要分支 / 循环 / 并行的复杂流程",
            "Send + State Reducer 是并行编排的核心",
            "Subgraph 把循环逻辑封装，主图保持简洁",
            "Middleware 解决合规与成本两大横切问题",
            "with_structured_output 让 LLM 输出强类型可靠",
            "LangSmith 三件套（Trace + Evals + Platform）补齐工程化",
            "Context Engineering 比 Prompt Engineering 更上层",
        ], font_size=12)

        _add_text_box(s, 17.5, 3, 15, 0.8, "下周方向", font_size=15, bold=True, color=ACCENT)
        _add_bullets(s, 17.5, 4.2, 15, 6, [
            "Skills 模块化与 anthropics/skills 仓库",
            "MCP（Model Context Protocol）协议",
            "Tools 规模化：把质检 KB 包成 MCP server",
            "Agent SDK 与 Subagent 调度机制",
            "可能的实验：把 KB 检索从本地 JSON 升级成向量库",
        ], font_size=12)

        _add_text_box(s, 1.5, 14, 30, 1.5,
                      "项目仓库：github.com/wensl121/ai-call-quality",
                      font_size=14, bold=True, color=DARK)
        _add_text_box(s, 1.5, 15.2, 30, 1,
                      "8 个 commit，2000+ 行代码 / prompt / 测试，11 个单测全过",
                      font_size=12, color=GREY)
        return s

    builders = [slide_1, slide_2, slide_3, slide_4, slide_5, slide_6,
                slide_7, slide_8, slide_9, slide_10, slide_11, slide_12,
                slide_13, slide_14, slide_15, slide_16, slide_17, slide_18, slide_19]

    for i, fn in enumerate(builders):
        slide = fn()
        if i > 0:  # 封面不要页码
            _add_footer(slide, i + 1, len(builders))

    prs.save(str(OUT))
    print(f"[done] saved: {OUT}  ({len(builders)} slides)")


if __name__ == "__main__":
    build()
