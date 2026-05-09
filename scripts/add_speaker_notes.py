"""给周分享2.pptx 每页加上演讲者备注（speaker notes）。

用法：
    python scripts/add_speaker_notes.py

会读取 C:\\Users\\Z\\Desktop\\周分享\\周分享2.pptx，
为每页添加 notes_text_frame，再保存到同一文件（PowerPoint 要先关闭）。
也可以通过 PPT_OUT 环境变量指定输出路径。
"""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation

IN = Path(r"C:\Users\Z\Desktop\周分享\周分享2.pptx")
OUT = Path(os.environ.get("PPT_OUT", str(IN)))

# ============================================================
# 演讲者备注（按 slide 1..22 顺序）
# 每条都按口语化风格写，方便讲的时候直接念
# ============================================================
NOTES: list[str] = [
    # 1. 封面
    """各位同事、领导大家好。我是开发三组的温舒麟，今天是第二周的预研分享。

上周我们看了 Claude Code、Hermes Agent、WorkBuddy 这些产品形态，今天我们换一个视角，从框架本身入手——重点讲 LangChain 1.0 和 LangGraph 的编排能力，以及上下文工程这个新趋势。

整个分享大概 25 分钟，中间会用一个我这周做的金融通话质检项目当例子来印证概念，最后留几分钟讨论。""",

    # 2. 目录
    """整个分享分三个模块。

第一个模块讲 LangChain Agent 架构——也就是 LangChain 1.0 里 create_agent 这套单 Agent 标准模式，以及它的 Middleware 横切机制。

第二个模块讲 LangGraph 图编排——这是用来处理复杂工作流的，包括 Graph Runtime、State 和 Reducer、Send 并行、还有 Subgraph 子图这几个核心概念。

第三个模块讲 Context Engineering，就是上下文工程，再加上 LangSmith 这套工程化的工具链——Tracing、Evals、Studio 和 Platform 部署。

每个模块讲完都会有一两页项目案例，把概念和真实代码对上。""",

    # 3. 01-1 Agent 编排概览
    """先说背景。LangChain 1.0 这次发版做了一个很关键的产品决策：把"单 Agent"和"复杂工作流"明确分到两个产品里。

左边的 create_agent 适合做单一职责的 Agent，比如客服 chatbot、tool-using agent、或者经典的 ReAct 模式。它的特点是简单——一个模型加一个工具循环就行，控制粒度中等。

右边的 LangGraph 适合复杂工作流，比如审批流、质检系统、研究流水线、或者人机协同的场景。它的控制粒度更高——每条边、每个节点都能定制，状态管理也更强。

中间通过 Middleware 拼接共性能力，这个我们后面会专门讲一页。

我这周做的质检系统因为有循环、有并行、有审核回退，所以选了 LangGraph。""",

    # 4. 01-2 create_agent
    """先看 create_agent，这是最简单的一种用法。

代码在左边，可以看到几乎是一行起一个 Agent。定义一个工具函数 get_weather，加上 @tool 装饰器，然后调 create_agent 传进去模型名、工具列表、和 system prompt 就完了。

它内部跑的就是经典的 ReAct loop——用户问问题，模型判断要不要调工具；要的话生成 tool call、执行、把结果回填给模型；不要就直接回答。这个循环一直跑到模型说"完成"为止。

这个模式很好用，但有局限——它只能处理"单个 Agent + 工具循环"。一旦你想加分支、加循环、加并行，create_agent 就力不从心了，这时候就要切到 LangGraph。""",

    # 5. 01-3 Middleware
    """说到 Middleware，LangChain 1.0 提供了三个标准 hook 点。

before_model 是调模型之前——可以注入上下文、做 PII 脱敏、做权限检查、动态生成 system prompt。

after_model 是调模型之后——日志审计、成本统计、输出审查、格式校验都放这里。

wrap_model_call 是完全包装——重试降级、响应缓存、模型故障切换、限流，这种全包式的逻辑就用它。

为什么要把这些抽出来？因为 PII 脱敏、成本统计、审计这些功能跟具体业务无关，每个节点都需要——抽成 middleware 之后，业务节点零侵入，将来要换实现也只改一个地方。

而且 middleware 可以叠加，像洋葱一样层层包装，每一层做一件事。""",

    # 6. 01 项目案例：PII + 成本追踪
    """讲完 Middleware 概念，看一下我项目里怎么用的。

左上是 PII 脱敏。我写了一个 redact_pii 函数，在 input_node 里调一次，整个下游所有节点——包括所有 LLM 调用、最终的 evidence 引用——都只会看到脱敏后的版本。手机号变成 PHONE 占位符、身份证变成 ID_CARD、地址变成 ADDRESS。对外输出的报告天然合规。

左下是成本追踪。我自己写了一个 invoke_structured 的 helper，用 include_raw=True 拿到 AIMessage 的 usage_metadata，把每次 LLM 调用的 token 数、估算成本记录下来。

右边是实测结果。一通 12 秒的通话总成本是 0.0024 美元——0.24 美分。rule_scorer 占了 67%。还能看到一个有意思的现象：question_extractor 调用了 2 次——这是 LangChain 结构化输出的隐藏 retry。这种以前根本看不见的成本，加了追踪后立刻浮出来。""",

    # 7. 01 with_structured_output
    """再看一个跟 Middleware 性质类似的工程化能力——结构化输出。

痛点是这样：以前我们让 LLM 返回 JSON，就是 prompt 里写"请返回 JSON 格式"，然后回来 json.loads 一下。这种做法有两个问题：第一，模型偶尔会加 markdown 围栏或者漏字段，json.loads 直接崩；第二，模型经常偷懒——比如让它返回 hot_words，它就给你个空数组糊弄。

LangChain 提供了 with_structured_output——左上定义一个 Pydantic schema，右上就能直接 invoke_structured 调用。框架自己负责 schema validate 加自动 retry。

下面这个表是同一段对话两种方式的对比。raw JSON 解析的时候，hot_words、business_words、should_say 全是空数组——模型在偷懒。换成 structured output 之后，每个字段都按 schema 强制填齐，hot_words 出了 5 个高价值词。""",

    # 8. 02-1 为什么需要图编排
    """进入第二个模块——LangGraph。

为什么单 Agent 不够？因为有三类问题它解决不了。

第一是分支——根据中间结果选不同后续路径。比如我项目里命中合规风险就直接判 0 分，非致命就继续累加扣分。这种条件分支 ReAct loop 写起来非常别扭。

第二是循环——重复执行直到满足条件。我项目里"打分→审核→不通过回去重打分"，最多跑 3 次。

第三是并行——多任务同时跑、最后合并。比如 30 条规则要各自评分，串行跑 150 秒，并行跑 15 秒。这是数量级的差距。

LangGraph 的解法是把流程画成有向图——节点是函数、边是流转关系，框架自己管 state、调度、并发、循环。""",

    # 9. 02-2 Graph 核心概念
    """LangGraph 的核心就四个组件。

Node 节点——一个纯函数，输入 state、返回 state 增量。可以纯计算，也可以调 LLM 或工具。

Edge 边——节点之间的连接关系，分三种：固定边、条件边、还有 Send 动态边。

State 状态——用 TypedDict 定义结构，加上 Reducer 决定多个写入怎么合并。这是 LangGraph 跟其他框架最大的区别——状态是一等公民。

Subgraph 子图——把一组节点封成单元，作为父图的"一个节点"用。

下面是最小可运行例子。定义一个 State，写一个 greet 节点，加边、compile、invoke——10 行代码起一个图。

执行模型简单说就是按 BFS 推进——当前层所有节点完成才触发下一层。compile 之后是不可变对象，多次 invoke 共享同一个 graph，性能很好。""",

    # 10. 02-3 State Reducer
    """重点讲一下 State Reducer，这是并行写入的关键。

默认行为是覆盖——左边代码里 deductions 字段没用 Annotated。如果有 5 个并行 scorer 同时写它，只有最后一个写入留下，前 4 个被覆盖——数据就丢了。

正确做法在右边——用 Annotated 加一个 reducer 函数。我自定义了一个 list_add_or_reset：右边传 RESET 信号就清空，否则就 append。

这个 RESET 信号在循环场景特别有用。我项目里 scoring_dispatcher 进入新一轮就先发 RESET 把 deductions 清空，然后 5 个 rule_scorer 各自 append，auditor 看到完整列表做判断。审核不通过回到 dispatcher，再次 RESET 开始下一轮，不污染历史数据。

LangGraph 内置一个 add_messages reducer 给消息列表用，自动追加。其余字段都用自定义 reducer。""",

    # 11. 02-4 Send + Reducer 综合应用
    """这一页把 Send 和 Reducer 放一起讲，因为它们是并行编排的最小可用组合。

左边是核心代码——scoring_dispatcher 节点先发 RESET 清空，然后 fan_out 函数为每条规则生成一个 Send 实例。Send 第一个参数是目标节点名，第二个是 payload。比如 5 条规则就生成 5 个 Send，加上一个 extraction_node 的 Send，总共 6 个并行实例。

所有 Send 完成后自动 fan-in 到 auditor 节点。

右边性能表是我跑出来的实测数据——这里我之前用了 25 秒做对比基准，是用 max_concurrency=1 强制串行测出来的。

注意一点：1 条规则的时候并行没优势，反而多了 fan-out 开销；5 条规则就开始有 50% 的收益；30 条规则估算从 150 秒降到 13-15 秒，差不多一个数量级。

而且——并行只压延迟，token 数和正确性都不变。这是免费的午餐。""",

    # 12. 02-5 Send 机制详细
    """这一页深入看 Send 的机制。

代码在左边——条件边的回调函数返回的不是节点名字符串，而是一个 Send 列表。每个 Send 携带专属 payload，调度时启动对应节点的一个新实例。

注意右下：实例可以是同名节点的多个并行执行——5 个 Send 都指向 rule_scorer，就会启动 5 个 rule_scorer 实例同时跑。

下游边自动等待所有 Send 完成才触发——这就是 fan-in。框架帮你做了同步。

这个机制很通用——经典用法是 map-reduce，每条数据派一个 worker。我们的并行规则评估就是这种。还有多查询并发检索、多模型对比评估等等都能用。""",

    # 13. 02-6 Subgraph
    """讲完 Send 讲 Subgraph——子图。

左边是主图，5 个节点：input、question_extractor、knowledge_retriever、scoring_with_audit、aggregator。注意 scoring_with_audit 这个节点不是普通函数——它是 build_scoring_subgraph 编译出来的子图，作为一个节点用。

右边是子图本身：4 个节点加一个内部循环。dispatcher 派 Send，rule_scorer 和 extraction_node 并行，最后 auditor 判断要不要重打分。

为什么要抽 Subgraph？五个收益：

主图清晰——从外面看就是一条流水线，循环细节封在里面。
可独立测试——子图可以单独 compile 加单测。
可复用——同样的"draft + critique"模式可以用在别处。
Studio 可视化——子图节点可以双击展开。
状态自动透传——父子图共享 GraphState，零 plumbing 代码。

这是模块化的基本功，写复杂图必须掌握。""",

    # 14. 02-6 项目案例：架构演进
    """这一页是项目实战的可视化对比。

左边是 v1，最早的版本——6 个节点线性排列，scoring_loop 和 auditor 之间有一条反馈边。从图上能看出来反馈边把整个图搞得有点乱。

右边是 v2，当前版本——主图变成 5 节点流水线，scoring_with_audit 是一个紫色框的子图。Studio 里可以双击展开看内部，内部就是 dispatcher、并行的 rule_scorer 和 extraction_node、最后 auditor。

演进路径有三步——先把 scoring_loop 和 auditor 抽成 subgraph，主图变干净；然后把 scoring_loop 拆成 dispatcher 加多个 rule_scorer 和 extraction，用 Send 并行；最后 Reducer 把并行结果合并起来。

延迟数据：5 规则评分串行 25 秒、并行 12.5 秒，30 规则估算 13-15 秒。""",

    # 15. 03-1 Prompt → Context Engineering
    """进入第三个模块——Context Engineering。

这是个新概念，但其实很多人已经在做了。表格对比说得很清楚：

Prompt Engineering 关注"单条 prompt 怎么写好"——措辞、few-shot 例子、角色扮演——产物是一条好 prompt。

Context Engineering 关注"上下文系统怎么设计"——动态检索、摘要压缩、工具描述、schema 注入——产物是一个能持续供给好上下文的系统。

可改性也不一样——prompt 改字符串就行，context 要动架构：加新节点、加新检索器、改 state 结构。

下面 Karpathy 那句话很经典——"软件 3.0 不是写代码，是工程化地构造模型上下文"。这个观点我觉得很对——我们以后做 AI 应用，主要工作不是 prompt 调优，而是设计上下文系统。""",

    # 16. 03-2 上下文设计的几个层次
    """既然 Context Engineering 是设计上下文系统，那系统包括哪些层次？

第一层是结构化输出——刚才项目案例里讲过，用 Pydantic schema 约束 LLM 返回，hot_words 从空数组变成 5 个高价值词。这是最简单的"约束模型输出形态"。

第二层是知识检索——RAG 也好、KB lookup 也好，从向量库或关键词库取参考答案，注入进 prompt 给模型比对。我项目里客户问题查 KB，答非所问就扣 rule 17。

第三层是 Few-shot 示例——动态注入历史的好例子或坏例子。模型看过几个对的就更会了。下周可以做的实验：从 Evals 失败的例子里挑反例注入进去。

第四层是反馈循环——上一轮的输出回馈给下一轮的输入。我项目里的审核机制就是这种——审核给出 issues，下一轮 scoring 收到 issues 做修正。

这四层从下到上，复杂度递增、但效果叠加。""",

    # 17. 03-3 LangSmith Tracing
    """LangSmith 是 LangChain 官方的工程化工具，先看 Tracing。

接入特别简单——env 里设 4 个变量，LANGCHAIN_TRACING_V2=true、API_KEY、PROJECT 名，LangChain 自动上传每次 LLM 调用，不用写一行追踪代码。

截图是我项目的 trace 视图。左边是 run 列表，每条对应一次调用。中间是 Waterfall——能看到 LangGraph 节点和内部 LLM 调用的层级关系，时间也很直观——哪个慢一目了然。右边是详情——input、output、token 用量都能看。

失败的 trace 会标红、可以重放。还支持搜索、过滤、按 tag 分组。

调试复杂图、定位 bug、做性能 profile，这是必备工具。""",

    # 18. 03-4 LangSmith Evals
    """光有 trace 不够——改 prompt 之后怎么知道效果是好了还是差了？这就要用 Evals。

左边截图是黄金数据集——我手工写了 4 个例子。第一个 clean_compliant 是合规通话，期望不扣分；missing_greeting 缺标准话术，期望命中 rule 7；fatal_guarantee 客服说了"保证收益"是致命违规，应该判 0 分；kb_mismatch 是答非所问，应该命中 rule 17。

右边截图是评估器结果。我写了 4 个评估器：rule_match 看 Jaccard 命中率、fatal_correctness 严格匹配致命触发、score_in_range 看最终分是否在区间内、review_match 看人工复核标记是否对。

baseline 跑出来：rule_match 75%、fatal_correctness 75%、另外两个 100%。

这 25% 的失配不是 bug，是信号——它告诉我"prompt 在这两个例子上跟期望对不上"。下次改 prompt 我就有定量基准——"rule_match 从 75% 到 92%、没有掉过去通过的例子"——这是用 Evals 之后才有的回归保险。""",

    # 19. 03-5 Studio + Platform 部署
    """LangSmith 的另两件套——Studio 和 Platform。

左边 Studio 是本地调试工具。装了 langgraph-cli 之后跑 langgraph dev，浏览器里就能打开 Studio。看到 graph 实时展开、可以单步执行、可以重放历史 run、还能直接编辑 input 重跑。

右边 Platform 是云端部署。我已经把项目部署上去了。流程很简单——langgraph.json 配好、推到 GitHub、LangSmith Platform 一键部署，最后给你一个 HTTPS API URL，可以直接 curl 或者用 SDK 调。环境变量在 Platform UI 配，DEEPSEEK_API_KEY 这种 secret 都有专门的安全存储。

部署之后还有 trace 自动记录、监控仪表板、版本管理——LangChain 把工程化全套都做好了。""",

    # 20. 03 彩蛋：Claude Code 协作
    """这一页是个彩蛋——其实整个项目和这份 PPT 都是用 Claude Code 协作完成的。

截图是我之前用 Claude Code 重构 audit 模块时候的终端，能看到它在写代码、跑测试、做 git commit。

具体感受是这样：从需求文档到项目骨架大概 30 分钟；每次增量重构——加结构化输出、加 Send 并行、抽 Subgraph——平均一小时左右；加 Middleware 那种横切的能力基本上零侵入；甚至 Evals 黄金数据集和评估器代码也是一起写出来的。

整个仓库 11 个 commit、2000 多行代码、11 个测试、Platform 部署完整可用——这种迭代速度以前是不可想象的。

下周第三周分享 Claude Agent SDK 和 MCP，会重点讲怎么自己用 Claude Code 这种工具。""",

    # 21. 03 总结 + 下周
    """最后做一个总结。

本周的关键认知：
create_agent 解决单 Agent 标准模式；复杂流程要靠 LangGraph 显式编排；Send、Reducer、Subgraph 是并行编排的三件套；Middleware 抽走横切关注点比如合规和成本；Context Engineering 比 Prompt Engineering 更上层；LangSmith 三件套——Trace、Evals、Platform——补齐工程化最后一公里；这些概念都已经在我项目里跑通验证过。

下周第三周预告：
讲 Skills 模块化——anthropics 那个 skills 仓库；MCP 协议；Tools 规模化——把我项目的 KB 包成一个 MCP server 让别的 Agent 也能用；Claude Agent SDK 和 Subagent 调度机制；可能的话做一个实验——把本地 KB 升级成向量库。

项目仓库地址 github.com/wensl121/ai-call-quality 已经公开，欢迎大家上去看代码。""",

    # 22. 结尾页
    """谢谢大家，分享到这里。有什么问题、想法、或者觉得哪里讲得不清楚，都欢迎现在或者会后找我聊。

仓库地址再贴一下：github.com/wensl121/ai-call-quality。下周见。""",
]


def add_notes(prs):
    n_slides = len(prs.slides)
    n_notes = len(NOTES)
    if n_slides != n_notes:
        print(f"⚠️ slide count {n_slides} != notes count {n_notes}; "
              f"will write min({n_slides},{n_notes}) and skip the rest")
    for i, slide in enumerate(prs.slides):
        if i >= len(NOTES):
            break
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = NOTES[i]


def main():
    if not IN.exists():
        raise SystemExit(f"input not found: {IN}")
    prs = Presentation(str(IN))
    add_notes(prs)
    prs.save(str(OUT))
    print(f"[done] saved: {OUT}  ({len(prs.slides)} slides, {len(NOTES)} notes written)")


if __name__ == "__main__":
    main()
